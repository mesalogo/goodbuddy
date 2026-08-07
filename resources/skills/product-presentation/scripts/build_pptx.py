#!/usr/bin/env python3
"""Build a 16:9 product presentation from a reviewed deck JSON file."""

import argparse
import json
import re
import tempfile
from pathlib import Path

from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


DEFAULT_THEME = {
    "primary": "1F3864",
    "accent": "C00000",
    "background": "FFFFFF",
    "text": "1A1A1A",
    "muted": "666666",
    "light": "F2F5F9",
    "font_zh": "Microsoft YaHei",
    "font_en": "Arial",
}
SLIDE_TYPES = {
    "title",
    "section",
    "bullets",
    "two-column",
    "metrics",
    "image",
    "closing",
}
NOTES_OPTIONAL_TYPES = {"section"}
MAX_IMAGE_BYTES = 40 * 1024 * 1024
MAX_IMAGE_PIXELS = 80_000_000
RENDER_DPI = 150


def color(value):
    value = value.lstrip("#")
    if not re.fullmatch(r"[0-9a-fA-F]{6}", value):
        raise ValueError(f"颜色必须是六位十六进制值：{value!r}")
    return RGBColor.from_string(value.upper())


def nonempty(value, path):
    if not isinstance(value, str) or not value.strip():
        raise ValueError(f"{path} 必须是非空字符串")
    return value.strip()


def load_deck(path):
    input_path = Path(path).expanduser().resolve()
    with input_path.open(encoding="utf-8") as handle:
        deck = json.load(handle)
    if not isinstance(deck, dict):
        raise ValueError("deck 根节点必须是对象")
    nonempty(deck.get("title"), "title")
    slides = deck.get("slides")
    if not isinstance(slides, list) or not slides:
        raise ValueError("slides 必须是非空数组")
    theme = {**DEFAULT_THEME, **deck.get("theme", {})}
    for key in ("primary", "accent", "background", "text", "muted", "light"):
        color(theme[key])
    for index, slide in enumerate(slides):
        validate_slide(slide, index)
    deck["theme"] = theme
    return input_path, deck


def validate_bullets(values, path, maximum=6):
    if not isinstance(values, list) or not values:
        raise ValueError(f"{path} 必须是非空数组")
    if len(values) > maximum:
        raise ValueError(f"{path} 最多 {maximum} 项")
    for index, value in enumerate(values):
        text = nonempty(value, f"{path}[{index}]")
        if len(text) > 120:
            raise ValueError(f"{path}[{index}] 超过 120 字")


def validate_slide(slide, index):
    path = f"slides[{index}]"
    if not isinstance(slide, dict):
        raise ValueError(f"{path} 必须是对象")
    slide_type = slide.get("type")
    if slide_type not in SLIDE_TYPES:
        raise ValueError(f"{path}.type 未知：{slide_type!r}")
    nonempty(slide.get("title"), f"{path}.title")
    if slide_type in {"bullets", "closing"}:
        validate_bullets(slide.get("bullets"), f"{path}.bullets")
    elif slide_type == "two-column":
        for side in ("left", "right"):
            column = slide.get(side)
            if not isinstance(column, dict):
                raise ValueError(f"{path}.{side} 必须是对象")
            nonempty(column.get("title"), f"{path}.{side}.title")
            validate_bullets(
                column.get("bullets"),
                f"{path}.{side}.bullets",
                maximum=5,
            )
    elif slide_type == "metrics":
        metrics = slide.get("metrics")
        if not isinstance(metrics, list) or not 1 <= len(metrics) <= 4:
            raise ValueError(f"{path}.metrics 必须包含 1 至 4 项")
        for metric_index, metric in enumerate(metrics):
            if not isinstance(metric, dict):
                raise ValueError(
                    f"{path}.metrics[{metric_index}] 必须是对象"
                )
            for field in ("value", "label", "detail"):
                nonempty(
                    metric.get(field),
                    f"{path}.metrics[{metric_index}].{field}",
                )
    elif slide_type == "image":
        nonempty(slide.get("image"), f"{path}.image")
    if slide_type not in NOTES_OPTIONAL_TYPES:
        nonempty(slide.get("notes"), f"{path}.notes")


def add_run_font(run, theme, size, bold=False, color_value=None):
    run.font.name = theme["font_en"]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color(color_value or theme["text"])
    run.font._element.set("lang", "zh-CN")


def add_textbox(
    slide,
    theme,
    x,
    y,
    width,
    height,
    text="",
    size=24,
    bold=False,
    color_value=None,
    align=PP_ALIGN.LEFT,
    vertical=MSO_ANCHOR.TOP,
    margin=0.08,
):
    shape = slide.shapes.add_textbox(
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height),
    )
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = vertical
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    add_run_font(run, theme, size, bold=bold, color_value=color_value)
    return shape


def set_background(slide, theme, key="background"):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color(theme[key])


def add_header(slide, theme, title):
    add_textbox(
        slide,
        theme,
        0.65,
        0.35,
        11.9,
        0.65,
        title,
        size=26,
        bold=True,
        color_value=theme["primary"],
    )
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.65),
        Inches(1.07),
        Inches(1.15),
        Inches(0.06),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color(theme["accent"])
    line.line.fill.background()


def add_footer(slide, deck, number):
    footer = deck.get("footer", "")
    theme = deck["theme"]
    if footer:
        add_textbox(
            slide,
            theme,
            0.65,
            7.08,
            10.8,
            0.22,
            footer,
            size=8.5,
            color_value=theme["muted"],
        )
    add_textbox(
        slide,
        theme,
        12.0,
        7.02,
        0.55,
        0.25,
        str(number),
        size=9,
        color_value=theme["muted"],
        align=PP_ALIGN.RIGHT,
    )


def add_bullet_frame(slide, theme, x, y, width, height, bullets, size=21):
    shape = slide.shapes.add_textbox(
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height),
    )
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.08)
    for index, item in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.space_after = Pt(12)
        paragraph.line_spacing = 1.12
        paragraph.font.size = Pt(size)
        paragraph.font.name = theme["font_en"]
        paragraph.font.color.rgb = color(theme["text"])
    return shape


def add_notes(slide, notes):
    if not notes:
        return
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = notes


def render_title(slide, deck, item):
    theme = deck["theme"]
    set_background(slide, theme, "primary")
    add_textbox(
        slide,
        theme,
        0.9,
        1.75,
        11.5,
        1.3,
        item["title"],
        size=34,
        bold=True,
        color_value="FFFFFF",
        vertical=MSO_ANCHOR.MIDDLE,
    )
    subtitle = item.get("subtitle", deck.get("subtitle", ""))
    if subtitle:
        add_textbox(
            slide,
            theme,
            0.95,
            3.2,
            10.8,
            0.8,
            subtitle,
            size=20,
            color_value="DCE6F1",
        )
    author = deck.get("author", "")
    if author:
        add_textbox(
            slide,
            theme,
            0.95,
            6.35,
            10.0,
            0.35,
            author,
            size=12,
            color_value="DCE6F1",
        )


def render_section(slide, deck, item):
    theme = deck["theme"]
    set_background(slide, theme, "light")
    add_textbox(
        slide,
        theme,
        1.0,
        2.2,
        11.2,
        1.0,
        item["title"],
        size=32,
        bold=True,
        color_value=theme["primary"],
        align=PP_ALIGN.CENTER,
        vertical=MSO_ANCHOR.MIDDLE,
    )
    if item.get("subtitle"):
        add_textbox(
            slide,
            theme,
            1.5,
            3.35,
            10.2,
            0.7,
            item["subtitle"],
            size=18,
            color_value=theme["muted"],
            align=PP_ALIGN.CENTER,
        )


def render_bullets(slide, deck, item):
    theme = deck["theme"]
    set_background(slide, theme)
    add_header(slide, theme, item["title"])
    add_bullet_frame(slide, theme, 0.9, 1.45, 11.5, 5.2, item["bullets"])


def render_two_column(slide, deck, item):
    theme = deck["theme"]
    set_background(slide, theme)
    add_header(slide, theme, item["title"])
    for x, column in ((0.75, item["left"]), (6.78, item["right"])):
        panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(1.55),
            Inches(5.55),
            Inches(4.95),
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = color(theme["light"])
        panel.line.color.rgb = color("D8E0EA")
        add_textbox(
            slide,
            theme,
            x + 0.3,
            1.8,
            4.95,
            0.5,
            column["title"],
            size=20,
            bold=True,
            color_value=theme["primary"],
        )
        add_bullet_frame(
            slide,
            theme,
            x + 0.25,
            2.5,
            5.0,
            3.55,
            column["bullets"],
            size=17,
        )


def render_metrics(slide, deck, item):
    theme = deck["theme"]
    set_background(slide, theme)
    add_header(slide, theme, item["title"])
    metrics = item["metrics"]
    gap = 0.25
    total_width = 11.8
    width = (total_width - gap * (len(metrics) - 1)) / len(metrics)
    for index, metric in enumerate(metrics):
        x = 0.75 + index * (width + gap)
        panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(1.75),
            Inches(width),
            Inches(4.45),
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = color(theme["light"])
        panel.line.color.rgb = color("D8E0EA")
        add_textbox(
            slide,
            theme,
            x + 0.15,
            2.05,
            width - 0.3,
            1.0,
            metric["value"],
            size=29,
            bold=True,
            color_value=theme["accent"],
            align=PP_ALIGN.CENTER,
            vertical=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide,
            theme,
            x + 0.15,
            3.15,
            width - 0.3,
            0.6,
            metric["label"],
            size=17,
            bold=True,
            color_value=theme["primary"],
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            theme,
            x + 0.2,
            4.0,
            width - 0.4,
            1.45,
            metric["detail"],
            size=12,
            color_value=theme["muted"],
            align=PP_ALIGN.CENTER,
        )


def prepare_image(image_path, max_width_in, max_height_in, staging_dir):
    """Normalize orientation and cap pixels so decks stay a usable size."""
    if image_path.stat().st_size > MAX_IMAGE_BYTES:
        raise ValueError(
            f"图片超过 {MAX_IMAGE_BYTES // (1024 * 1024)}MB：{image_path.name}"
        )
    with Image.open(image_path) as image:
        if image.width * image.height > MAX_IMAGE_PIXELS:
            raise ValueError(f"图片像素数过大：{image_path.name}")
        image = ImageOps.exif_transpose(image)
        width, height = image.size
        scale = min(max_width_in / width, max_height_in / height)
        target = (
            max(1, round(width * scale * RENDER_DPI)),
            max(1, round(height * scale * RENDER_DPI)),
        )
        if target[0] >= width and target[1] >= height:
            if image_path.suffix.lower() in {".png", ".jpg", ".jpeg", ".gif"}:
                return image_path, width, height
            target = (width, height)
        resized = image.resize(target, Image.LANCZOS)
        if resized.mode not in ("RGB", "RGBA", "L"):
            resized = resized.convert("RGBA")
        staged = staging_dir / f"{image_path.stem}-{target[0]}x{target[1]}.png"
        resized.save(staged, format="PNG", optimize=True)
    return staged, width, height


def render_image(slide, deck, item, base_dir, staging_dir):
    theme = deck["theme"]
    set_background(slide, theme)
    add_header(slide, theme, item["title"])
    image_path = (base_dir / item["image"]).resolve()
    if base_dir != image_path and base_dir not in image_path.parents:
        raise ValueError(f"图片必须位于 deck.json 目录内：{image_path}")
    if not image_path.is_file():
        raise FileNotFoundError(f"图片不存在：{image_path}")
    max_width, max_height = 11.5, 5.35
    source, width, height = prepare_image(
        image_path, max_width, max_height, staging_dir
    )
    scale = min(max_width / width, max_height / height)
    draw_width, draw_height = width * scale, height * scale
    slide.shapes.add_picture(
        str(source),
        Inches((13.333 - draw_width) / 2),
        Inches(1.35 + (5.35 - draw_height) / 2),
        width=Inches(draw_width),
        height=Inches(draw_height),
    )
    if item.get("caption"):
        add_textbox(
            slide,
            theme,
            1.0,
            6.55,
            11.3,
            0.3,
            item["caption"],
            size=10,
            color_value=theme["muted"],
            align=PP_ALIGN.CENTER,
        )


def render_closing(slide, deck, item):
    theme = deck["theme"]
    set_background(slide, theme, "primary")
    add_textbox(
        slide,
        theme,
        0.9,
        1.2,
        11.5,
        0.9,
        item["title"],
        size=32,
        bold=True,
        color_value="FFFFFF",
        align=PP_ALIGN.CENTER,
    )
    add_bullet_frame(
        slide,
        {**theme, "text": "FFFFFF"},
        2.0,
        2.55,
        9.3,
        3.2,
        item["bullets"],
        size=20,
    )


RENDERERS = {
    "title": render_title,
    "section": render_section,
    "bullets": render_bullets,
    "two-column": render_two_column,
    "metrics": render_metrics,
    "closing": render_closing,
}


def build_presentation(input_path, deck, output_path):
    output = Path(output_path).expanduser().resolve()
    if output.suffix.lower() != ".pptx":
        raise ValueError("输出必须使用 .pptx 扩展名")
    if output == input_path:
        raise ValueError("输出不能覆盖 deck.json")
    output.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]
    with tempfile.TemporaryDirectory(prefix="deck-images-") as staging:
        staging_dir = Path(staging)
        for index, item in enumerate(deck["slides"], 1):
            slide = presentation.slides.add_slide(blank)
            if item["type"] == "image":
                render_image(slide, deck, item, input_path.parent, staging_dir)
            else:
                RENDERERS[item["type"]](slide, deck, item)
            if item["type"] != "title":
                add_footer(slide, deck, index)
            add_notes(slide, item.get("notes", ""))
        presentation.save(output)
    return output


def parse_args():
    parser = argparse.ArgumentParser(description="从 deck.json 生成产品介绍 PPTX")
    parser.add_argument("--input", required=True, help="deck.json")
    parser.add_argument("--output", required=True, help="输出 .pptx")
    return parser.parse_args()


def main():
    args = parse_args()
    input_path, deck = load_deck(args.input)
    output = build_presentation(input_path, deck, args.output)
    print(f"saved: {output}")


if __name__ == "__main__":
    main()
