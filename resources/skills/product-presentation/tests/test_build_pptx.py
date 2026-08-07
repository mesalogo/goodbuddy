import copy
import importlib.util
import io
import json
import tempfile
import unittest
from pathlib import Path

from PIL import Image
from pptx import Presentation


SKILL_DIR = Path(__file__).resolve().parents[1]
SCRIPT = SKILL_DIR / "scripts" / "build_pptx.py"
SPEC = importlib.util.spec_from_file_location("build_pptx", SCRIPT)
build_pptx = importlib.util.module_from_spec(SPEC)
SPEC.loader.exec_module(build_pptx)


class BuildPptxTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        template = SKILL_DIR / "templates" / "deck.example.json"
        cls.deck = json.loads(template.read_text(encoding="utf-8"))

    def test_builds_example_deck(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            input_path = root / "deck.json"
            input_path.write_text(
                json.dumps(self.deck, ensure_ascii=False),
                encoding="utf-8",
            )
            loaded_path, deck = build_pptx.load_deck(input_path)
            output = build_pptx.build_presentation(
                loaded_path,
                deck,
                root / "deck.pptx",
            )
            presentation = Presentation(output)
            self.assertEqual(len(presentation.slides), len(self.deck["slides"]))
            all_text = "\n".join(
                shape.text
                for slide in presentation.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )
            self.assertIn("目标场景与问题", all_text)
            self.assertIn("下一步", all_text)
            for slide, item in zip(presentation.slides, self.deck["slides"]):
                self.assertEqual(
                    slide.notes_slide.notes_text_frame.text.strip(),
                    item["notes"].strip(),
                )

    def test_requires_notes_outside_section_slides(self):
        deck = copy.deepcopy(self.deck)
        deck["slides"][1].pop("notes")
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "deck.json"
            path.write_text(json.dumps(deck, ensure_ascii=False), encoding="utf-8")
            with self.assertRaisesRegex(ValueError, r"slides\[1\]\.notes"):
                build_pptx.load_deck(path)

    def test_downsamples_large_image_slide(self):
        deck = copy.deepcopy(self.deck)
        deck["slides"].insert(
            3,
            {
                "type": "image",
                "title": "参考架构",
                "image": "architecture.png",
                "caption": "示例架构图",
                "notes": "说明组件边界。",
            },
        )
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            Image.new("RGB", (5000, 3000), "white").save(root / "architecture.png")
            path = root / "deck.json"
            path.write_text(json.dumps(deck, ensure_ascii=False), encoding="utf-8")
            loaded_path, loaded = build_pptx.load_deck(path)
            output = build_pptx.build_presentation(
                loaded_path, loaded, root / "deck.pptx"
            )
            presentation = Presentation(output)
            self.assertEqual(len(presentation.slides), len(deck["slides"]))
            pictures = [
                shape
                for slide in presentation.slides
                for shape in slide.shapes
                if shape.shape_type == 13
            ]
            self.assertEqual(len(pictures), 1)
            with Image.open(io.BytesIO(pictures[0].image.blob)) as embedded:
                self.assertLess(embedded.width, 5000)

    def test_rejects_image_outside_deck_directory(self):
        deck = copy.deepcopy(self.deck)
        deck["slides"].insert(
            1,
            {
                "type": "image",
                "title": "外部图片",
                "image": "../outside.png",
                "notes": "越界图片。",
            },
        )
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp) / "deck"
            root.mkdir()
            Image.new("RGB", (10, 10), "white").save(Path(tmp) / "outside.png")
            path = root / "deck.json"
            path.write_text(json.dumps(deck, ensure_ascii=False), encoding="utf-8")
            loaded_path, loaded = build_pptx.load_deck(path)
            with self.assertRaisesRegex(ValueError, "必须位于"):
                build_pptx.build_presentation(
                    loaded_path, loaded, root / "deck.pptx"
                )

    def test_rejects_too_many_bullets(self):
        deck = copy.deepcopy(self.deck)
        deck["slides"][1]["bullets"] = [f"项目 {index}" for index in range(7)]
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "deck.json"
            path.write_text(
                json.dumps(deck, ensure_ascii=False),
                encoding="utf-8",
            )
            with self.assertRaisesRegex(ValueError, "最多 6 项"):
                build_pptx.load_deck(path)

    def test_rejects_unknown_slide_type(self):
        deck = copy.deepcopy(self.deck)
        deck["slides"][0]["type"] = "unknown"
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "deck.json"
            path.write_text(
                json.dumps(deck, ensure_ascii=False),
                encoding="utf-8",
            )
            with self.assertRaisesRegex(ValueError, "未知"):
                build_pptx.load_deck(path)


if __name__ == "__main__":
    unittest.main()
